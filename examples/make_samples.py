"""Generate sample .docx / .pptx / image fixtures for demoing the pipeline.

Usage: python examples/make_samples.py
"""

from __future__ import annotations

import json
from pathlib import Path

HERE = Path(__file__).parent
SAMPLE_TXT = (HERE / "sample_contract.txt").read_text(encoding="utf-8")


def make_docx() -> None:
    import docx

    doc = docx.Document()
    doc.add_heading("技术服务合同", level=1)
    for line in SAMPLE_TXT.splitlines()[2:]:
        doc.add_paragraph(line)
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "开户银行"
    table.cell(0, 1).text = "招商银行上海张江支行"
    table.cell(1, 0).text = "账号"
    table.cell(1, 1).text = "6225880123456789012"
    out = HERE / "sample_contract.docx"
    doc.save(str(out))
    print(f"wrote {out}")


def make_pptx() -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "项目汇报 — 客户：北京云启智能科技有限公司"
    slide.placeholders[1].text = (
        "联系人：陈建国\n电话：13812345678\n邮箱：chenjg@yunqi-tech.example.com"
    )
    out = HERE / "sample_report.pptx"
    prs.save(str(out))
    print(f"wrote {out}")


def _cjk_font(size: int = 28):
    """PIL's default font has no CJK glyphs; try common system fonts first."""
    from PIL import ImageFont

    candidates = [
        "msyh.ttc", "simhei.ttf", "simsun.ttc",  # Windows
        "/System/Library/Fonts/PingFang.ttc",  # macOS
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",  # Linux
    ]
    for name in candidates:
        try:
            return ImageFont.truetype(name, size)
        except OSError:
            continue
    return ImageFont.load_default()


def make_image() -> None:
    from PIL import Image, ImageDraw

    img = Image.new("RGB", (800, 200), "white")
    draw = ImageDraw.Draw(img)
    font = _cjk_font()
    line1 = "联系人：陈建国  电话：13812345678"
    line2 = "地址：北京市朝阳区望京东路8号"
    draw.text((20, 30), line1, fill="black", font=font)
    draw.text((20, 100), line2, fill="black", font=font)
    out = HERE / "sample_scan.png"
    img.save(out)
    sidecar = [
        {"text": line1, "bbox": [20, 25, 760, 55], "confidence": 0.98},
        {"text": line2, "bbox": [20, 95, 760, 125], "confidence": 0.97},
    ]
    (HERE / "sample_scan.png.ocr.json").write_text(
        json.dumps(sidecar, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    print(f"wrote {out} (+ .ocr.json sidecar for mock OCR)")


if __name__ == "__main__":
    make_docx()
    make_pptx()
    make_image()
