"""File type dispatch: read text per format, write redacted output per format.

Heavy format logic stays in the libraries (python-docx / python-pptx / pypdf);
agent-driven runs use the .claude/skills/ docx/pptx/pdf skills instead of this
module. Keep this simple — demo-grade fidelity is the goal.
"""

from __future__ import annotations

from pathlib import Path

from .redactors import redact_image, redact_text
from .schemas import Entity, OCRBlock

TEXT_EXTS = {".txt", ".md"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".webp"}


def file_type_of(path: str | Path) -> str:
    ext = Path(path).suffix.lower()
    if ext in TEXT_EXTS:
        return "text"
    if ext == ".docx":
        return "docx"
    if ext == ".pptx":
        return "pptx"
    if ext == ".pdf":
        return "pdf"
    if ext in IMAGE_EXTS:
        return "image"
    raise ValueError(f"不支持的文件类型: {ext}")


# ---------------------------------------------------------------------------
# Read
# ---------------------------------------------------------------------------


def read_document(path: str | Path) -> tuple[str, dict]:
    """Return (full_text, parse_stats). Images return empty text (OCR path)."""
    path = Path(path)
    ftype = file_type_of(path)
    if ftype == "text":
        text = path.read_text(encoding="utf-8", errors="replace")
        return text, {"解析方式": "纯文本", "字符数": len(text)}
    if ftype == "docx":
        return _read_docx(path)
    if ftype == "pptx":
        return _read_pptx(path)
    if ftype == "pdf":
        return _read_pdf(path)
    return "", {"解析方式": "图片，待 OCR"}


def _read_docx(path: Path) -> tuple[str, dict]:
    import docx

    doc = docx.Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    n_tables = len(doc.tables)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    parts.append(cell.text)
    return "\n".join(parts), {
        "解析方式": "docx skill (python-docx)",
        "段落数": len(doc.paragraphs),
        "表格数": n_tables,
    }


def _read_pptx(path: Path) -> tuple[str, dict]:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    n_shapes = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                n_shapes += 1
                parts.append(shape.text_frame.text)
    return "\n".join(parts), {
        "解析方式": "pptx skill (python-pptx)",
        "幻灯片数": len(prs.slides._sldIdLst),
        "文本框数": n_shapes,
    }


def _read_pdf(path: Path) -> tuple[str, dict]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = [page.extract_text() or "" for page in reader.pages]
    text = "\n".join(pages)
    stats = {"解析方式": "pdf skill (pypdf)", "页数": len(pages), "字符数": len(text)}
    if len(text.strip()) < 20:
        stats["提示"] = "提取文本极少，可能是扫描件，需走 OCR 路径"
    return text, stats


# ---------------------------------------------------------------------------
# Write
# ---------------------------------------------------------------------------


def write_redacted(
    path: str | Path,
    text: str,
    entities: list[Entity],
    output_dir: str | Path,
    ocr_blocks: list[OCRBlock] | None = None,
) -> tuple[list[Path], list[str]]:
    """Write redacted output files; returns (paths, warnings)."""
    path = Path(path)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    ftype = file_type_of(path)
    stem = path.stem
    redacted = redact_text(text, entities)
    outputs: list[Path] = []
    warnings: list[str] = []

    if ftype == "text":
        out = output_dir / f"{stem}.redacted{path.suffix}"
        out.write_text(redacted, encoding="utf-8")
        outputs.append(out)
    elif ftype == "docx":
        outputs.append(_write_docx(path, entities, output_dir / f"{stem}.redacted.docx"))
        md = output_dir / f"{stem}.redacted.md"
        md.write_text(redacted, encoding="utf-8")
        outputs.append(md)
    elif ftype == "pptx":
        outputs.append(_write_pptx(path, entities, output_dir / f"{stem}.redacted.pptx"))
        md = output_dir / f"{stem}.redacted.md"
        md.write_text(redacted, encoding="utf-8")
        outputs.append(md)
    elif ftype == "pdf":
        out = output_dir / f"{stem}.redacted.md"
        out.write_text(redacted, encoding="utf-8")
        outputs.append(out)
        warnings.append("PDF 以提取文本方式输出 .redacted.md；扫描件请使用 OCR 路径")
    elif ftype == "image":
        out = output_dir / f"{stem}.redacted{path.suffix}"
        redact_image(path, entities, ocr_blocks or [], out)
        outputs.append(out)
    return outputs, warnings


def _replace_in_paragraph(paragraph, entities: list[Entity]) -> None:
    """Paragraph-level replace: rebuilds runs only when a hit crosses run text."""
    full = paragraph.text
    new = redact_text(full, entities)
    if new == full:
        return
    # try run-by-run first to preserve formatting
    for run in paragraph.runs:
        run.text = redact_text(run.text, entities)
    if paragraph.text == new:
        return
    # entity spans run boundaries: collapse into first run (demo-grade tradeoff)
    if paragraph.runs:
        paragraph.runs[0].text = new
        for run in paragraph.runs[1:]:
            run.text = ""


def _write_docx(src: Path, entities: list[Entity], out: Path) -> Path:
    import docx

    doc = docx.Document(str(src))
    for p in doc.paragraphs:
        _replace_in_paragraph(p, entities)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    _replace_in_paragraph(p, entities)
    doc.save(str(out))
    return out


def _write_pptx(src: Path, entities: list[Entity], out: Path) -> Path:
    from pptx import Presentation

    prs = Presentation(str(src))
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = redact_text(run.text, entities)
                # cross-run hits: collapse paragraph
                joined = "".join(r.text for r in para.runs)
                fixed = redact_text(joined, entities)
                if fixed != joined and para.runs:
                    para.runs[0].text = fixed
                    for r in para.runs[1:]:
                        r.text = ""
    prs.save(str(out))
    return out
